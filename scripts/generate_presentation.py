from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Credit Risk Model: Hybrid Explainability & Dashboard"
subtitle = slide.placeholders[1]
subtitle.text = "Interim Progress Report\nFeb 2026"

# Custom color theme
ACCENT = RGBColor(0, 102, 204)

# Section Slide


def add_section_slide(title_text):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT


# Content Slide


def add_content_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    body = slide.shapes.placeholders[1].text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = point
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT


# Add sections
add_section_slide("Plan vs. Progress")
add_content_slide(
    "Original Plan",
    [
        "Build credit risk pipeline with minimal leakage",
        "Temporal data split, robust processing",
        "Streamlit dashboard & API fallback",
        "SHAP explainability (global/local)",
        "Clear documentation & reproducibility",
    ],
)
add_content_slide(
    "Actual Progress",
    [
        "Temporal split & leakage reduction implemented",
        "Data pipeline & CLI completed",
        "Streamlit dashboard with batch/API scoring",
        "SHAP explainability (global, local, pie chart)",
        "MLflow tracking, updated README & demo",
    ],
)

add_section_slide("Completed Work")
add_content_slide(
    "Key Deliverables",
    [
        "Temporal split, proxy target, feature engineering",
        "MLflow model training, class_weight balancing",
        "FastAPI backend, Streamlit dashboard",
        "SHAP summary, bar, pie chart visuals",
        "Comprehensive README, demo GIF, code comments",
    ],
)

add_section_slide("Blockers & Challenges")
add_content_slide(
    "Challenges & Solutions",
    [
        "Initial leakage: fixed with temporal split",
        "SHAP input errors: fixed with validation",
        "Explainability for non-tech: added pie chart",
        "Linting/indentation: autopep8 & manual review",
    ],
)
add_content_slide(
    "What Remains",
    [
        "Full documentation refresh (CONTRIBUTING, API docs)",
        "Advanced monitoring/challenger analysis (if time)",
        "Prepare final report/codebase for submission",
    ],
)

add_section_slide("Demo & Visuals")

# Demo & Visuals Slide (with manual title)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
shapes = slide.shapes
# Add title textbox manually
title_box = shapes.add_textbox(Inches(1), Inches(0.3), Inches(7), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Dashboard Demo"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

img_path = "reports/figures/dashboard.gif"
try:
    shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(7))
except Exception:
    shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(7), Inches(2)).text = (
        "[Demo GIF not found]"
    )

# Save presentation
prs.save("reports/credit_risk_interim_report.pptx")
print("Presentation saved as reports/credit_risk_interim_report.pptx")
